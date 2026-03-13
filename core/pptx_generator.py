"""PowerPoint generator - replaces VBA macro logic for modifying PPTX template slides."""

import copy
import logging
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Emu

from core.models import ExtendedMarketData

logger = logging.getLogger(__name__)

# Colors matching VBA constants
COLOR_DOMINATING = RGBColor(23, 52, 97)       # Dark blue
COLOR_FASTEST_GROWING = RGBColor(117, 200, 146)  # Green
COLOR_TITLE_BLUE = RGBColor(23, 52, 97)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(0, 0, 0)

# VBA region-to-shape mapping for Slide 1
# Shape indices (0-based) for the 6 map region groups
REGION_SHAPE_MAP = {
    "Middle East": 0,
    "Latin America": 1,
    "Africa": 2,
    "North America": 3,
    "Europe": 4,
    "Asia Pacific": 5,
}


def _set_text(shape, text: str, font_size: int = None, bold: bool = None,
              color: RGBColor = None, alignment=None):
    """Set text content of a shape, preserving existing formatting where possible."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs:
        para = tf.paragraphs[0]
        # Clear existing runs
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = text
        else:
            run = para.add_run()
            run.text = text
        if font_size is not None:
            para.runs[0].font.size = Pt(font_size)
        if bold is not None:
            para.runs[0].font.bold = bold
        if color is not None:
            para.runs[0].font.color.rgb = color
        if alignment is not None:
            para.alignment = alignment


def _set_shape_fill_color(shape, color: RGBColor):
    """Set solid fill color on a shape or group."""
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    except Exception:
        # Groups may not support direct fill - try child shapes
        if hasattr(shape, 'shapes'):
            for child in shape.shapes:
                try:
                    child.fill.solid()
                    child.fill.fore_color.rgb = color
                except Exception:
                    pass


def _modify_regional_insights(slide, data: ExtendedMarketData):
    """Modify Slide 1 - Regional Insights.

    Reference shape mapping from Output.pptx:
    [0-5] Group shapes = map regions (Middle East, Latin America, Africa, North America, Europe, Asia Pacific)
    [8]   TextBox = "Regional Insights, {year}" title
    [10]  TextBox = "{region} - Estimated Market Revenue Share, {year}" subtitle
    [11]  TextBox = "{pct}%" percentage
    [12]  TextBox = market name
    [15]  TextBox = "Total Market Size: {value}"
    """
    shapes = slide.shapes

    # Color the map regions based on status
    for region in data.regions:
        idx = REGION_SHAPE_MAP.get(region.name)
        if idx is not None and idx < len(shapes):
            if region.status == "Dominating":
                _set_shape_fill_color(shapes[idx], COLOR_DOMINATING)
            elif region.status == "Fastest Growing":
                _set_shape_fill_color(shapes[idx], COLOR_FASTEST_GROWING)

    # Title: "Regional Insights, {year}"
    if len(shapes) > 8:
        _set_text(shapes[8], f"Regional Insights, {data.base_year + 1}")

    # Subtitle with dominating region
    dom_region = data.dominating_region or "N/A"
    if len(shapes) > 10:
        _set_text(shapes[10],
                  f"{dom_region} - Estimated Market Revenue Share, {data.base_year + 1}")

    # Percentage
    if len(shapes) > 11:
        pct_str = f"{data.market_share_pct * 100:.1f}%"
        _set_text(shapes[11], pct_str)

    # Market name
    if len(shapes) > 12:
        _set_text(shapes[12], data.market_name)

    # Total market size
    if len(shapes) > 15:
        _set_text(shapes[15], f"Total Market Size: {data.market_size_value}")


def _modify_impact_analysis(slide, data: ExtendedMarketData):
    """Modify Slide 2 - Impact Analysis of Key Factors.

    Reference shape mapping from Output.pptx:
    [2]  Rectangle = Driver 1 text
    [3]  Rectangle = Driver 2 text
    [4]  Rectangle = Restraint 1 text
    [5]  Rectangle = Restraint 2 text
    [6]  Rectangle = Opportunity 1 text
    [7]  Rectangle = Opportunity 2 text
    [11] TextBox = title
    [12-17] Groups = indicator arrows (position based on indicator value)
    """
    shapes = slide.shapes

    # Set driver/restraint/opportunity text
    labels = [
        (2, data.driver_1),
        (3, data.driver_2),
        (4, data.restraint_1),
        (5, data.restraint_2),
        (6, data.opportunity_1),
        (7, data.opportunity_2),
    ]
    for idx, text in labels:
        if idx < len(shapes):
            _set_text(shapes[idx], text)

    # Title
    if len(shapes) > 11:
        _set_text(shapes[11], f"Impact Analysis of Key Factors | {data.market_name}")

    # Position indicator arrows based on values
    # VBA formula: .Left = 340 + (100 * indicator_value) in points
    indicators = [
        (12, data.driver_1_indicator),
        (13, data.driver_2_indicator),
        (14, data.restraint_1_indicator),
        (15, data.restraint_2_indicator),
        (16, data.opportunity_1_indicator),
        (17, data.opportunity_2_indicator),
    ]
    for idx, value in indicators:
        if idx < len(shapes):
            try:
                shapes[idx].left = Pt(340 + (100 * value))
            except Exception as e:
                logger.debug(f"Could not position indicator shape[{idx}]: {e}")


def _modify_key_takeaways(slide, data: ExtendedMarketData):
    """Modify Slide 3 - Key Takeaways from Lead Analyst.

    Reference shape mapping from Output.pptx:
    [0] Rectangle = main takeaways text box
    [3] TextBox = "Key Takeaways from Lead Analyst" header
    """
    shapes = slide.shapes

    if not data.takeaways or len(shapes) < 1:
        return

    # Combine takeaways into single text block
    combined_text = "\n\n".join(data.takeaways)

    # Calculate font size based on text length (VBA formula)
    total_chars = sum(len(t) for t in data.takeaways)
    # Approximate lines: each line ~111 chars
    total_lines = sum((len(t) // 111) + 1 for t in data.takeaways)
    font_size = max(8, min(22, int(22 - 0.66 * (total_lines - 11.45))))

    shape = shapes[0]
    if shape.has_text_frame:
        tf = shape.text_frame
        # Clear existing content
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = ""

        # Set first paragraph with all takeaways
        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = combined_text
            tf.paragraphs[0].runs[0].font.size = Pt(font_size)
        else:
            para = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            run = para.add_run()
            run.text = combined_text
            run.font.size = Pt(font_size)


def _modify_segmental_insights(slide, data: ExtendedMarketData):
    """Modify Slide 4 - Segmental Insights with doughnut chart.

    Reference shape mapping from Output.pptx:
    [0]  TextBox = "{market_name}, By {segment_type}, {year}" title
    [2]  TextBox = "Total Market Size: {value}"
    [6]  TextBox = "{largest_segment} - Estimated Market Revenue Share, {year}" subtitle
    [7]  TextBox = "{pct}%" largest segment percentage
    [8]  TextBox = market name
    [13] Chart = doughnut chart
    """
    shapes = slide.shapes
    chart_data_items = data.segment_chart_data

    if not chart_data_items:
        return

    # Find largest segment
    largest = max(chart_data_items, key=lambda x: x.value)
    segment_type = data.segments[0].name if data.segments else "By Type"

    # Title
    if len(shapes) > 0:
        _set_text(shapes[0],
                  f"{data.market_name}, {segment_type}, {data.base_year + 1}")

    # Total market size
    if len(shapes) > 2:
        _set_text(shapes[2], f"Total Market Size: {data.market_size_value}")

    # Subtitle with largest segment
    if len(shapes) > 6:
        _set_text(shapes[6],
                  f"{largest.label} - Estimated Market Revenue Share, {data.base_year + 1}")

    # Largest segment percentage
    if len(shapes) > 7:
        _set_text(shapes[7], f"{largest.value * 100:.1f}%")

    # Market name
    if len(shapes) > 8:
        _set_text(shapes[8], data.market_name)

    # Update doughnut chart data
    if len(shapes) > 13 and shapes[13].has_chart:
        chart = shapes[13].chart
        cd = CategoryChartData()
        cd.categories = [item.label for item in chart_data_items]
        cd.add_series("Market Share", [item.value for item in chart_data_items])
        chart.replace_data(cd)

        # Format data labels
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.number_format = "0.0%"
            data_labels.font.size = Pt(8)
            data_labels.font.bold = True
            data_labels.font.color.rgb = COLOR_WHITE
        except Exception as e:
            logger.debug(f"Could not format chart data labels: {e}")

        # Set doughnut hole size
        try:
            chart.plots[0]._element.attrib.clear()
            # Access via XML to set hole size
            from pptx.oxml.ns import qn
            doughnut = chart.plots[0]._element
            hole_size = doughnut.find(qn("c:holeSize"))
            if hole_size is not None:
                hole_size.set("val", "50")
        except Exception:
            pass


def _modify_market_key_players(slide, data: ExtendedMarketData):
    """Modify Slide 5 - Market Key Players.

    Reference shape mapping from Output.pptx:
    [0]  TextBox = "Market Concentration, By Players" title
    [7]  Right Arrow = concentration indicator (position based on D31)
    [8]  Rectangle = market name
    [17+] Rectangles = company names (dynamically created in VBA)
    """
    shapes = slide.shapes

    # Position concentration arrow
    # VBA formula: .Top = 380 - (42 * concentration_value) in points
    if len(shapes) > 7:
        try:
            new_top = Pt(380 - (42 * data.concentration_value))
            shapes[7].top = new_top
        except Exception as e:
            logger.debug(f"Could not position arrow: {e}")

    # Market name
    if len(shapes) > 8:
        _set_text(shapes[8], data.market_name)

    # Update existing company name shapes (indices 17+)
    # The reference Output.pptx already has company rectangles at [17] onwards
    company_start_idx = 17
    for i, company in enumerate(data.companies):
        shape_idx = company_start_idx + i
        if shape_idx < len(shapes):
            _set_text(shapes[shape_idx], company)

    # If there are more existing shapes than companies, clear extras
    for i in range(len(data.companies), 20):
        shape_idx = company_start_idx + i
        if shape_idx < len(shapes) and shapes[shape_idx].has_text_frame:
            _set_text(shapes[shape_idx], "")


def generate_pptx(
    template_path: Path,
    data: ExtendedMarketData,
    output_path: Path,
) -> Path:
    """Generate modified PowerPoint from template using market data.

    Args:
        template_path: Path to the PPTX template file
        data: ExtendedMarketData with all fields populated
        output_path: Where to save the modified PPTX

    Returns:
        Path to the saved PPTX file
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(template_path))

    slides = list(prs.slides)
    if len(slides) < 5:
        logger.warning(f"Template has only {len(slides)} slides, expected 5")

    # Modify each slide
    slide_handlers = [
        (0, _modify_regional_insights, "Regional Insights"),
        (1, _modify_impact_analysis, "Impact Analysis"),
        (2, _modify_key_takeaways, "Key Takeaways"),
        (3, _modify_segmental_insights, "Segmental Insights"),
        (4, _modify_market_key_players, "Market Key Players"),
    ]

    for idx, handler, name in slide_handlers:
        if idx < len(slides):
            try:
                handler(slides[idx], data)
                logger.info(f"Modified slide {idx + 1}: {name}")
            except Exception as e:
                logger.error(f"Error modifying slide {idx + 1} ({name}): {e}")
        else:
            logger.warning(f"Skipping slide {idx + 1} ({name}): not in template")

    prs.save(str(output_path))
    logger.info(f"Saved PPTX to {output_path}")
    return output_path
